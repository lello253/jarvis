"""
file_processor.py — JARVIS Universal File Processor

Supported types:
  image   → describe, ocr, resize, convert, compress, crop
  pdf     → summarize, extract_text, extract_pages, to_word
  docx    → summarize, extract_text, reformat, translate_hint
  txt/md  → summarize, reformat, translate_hint, word_count
  csv     → analyze, filter, sort, convert, stats
  xlsx    → analyze, filter, convert, stats
  json    → validate, format, extract, convert
  code    → explain, review, fix, run, document
  audio   → transcribe, trim, convert, info
  video   → trim, extract_audio, extract_frame, info, compress
  zip     → list, extract
  pptx    → summarize, extract_text, to_pdf
"""

import os
import re
import json
import shutil
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime

def _get_api_key() -> str:
    config_path = Path(__file__).resolve().parent.parent / "config" / "api_keys.json"
    try:
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)["gemini_api_key"]
    except Exception:
        return ""


def _gemini_client():
    from google import genai
    _c = genai.Client(api_key=_get_api_key())

    class _W:
        def generate_content(self, contents):
            return _c.models.generate_content(model="gemini-2.5-flash", contents=contents)

    return _W()


def _detect_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    image_exts = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "svg", "ico"}
    video_exts = {"mp4", "avi", "mov", "mkv", "wmv", "flv", "webm", "m4v", "3gp"}
    audio_exts = {"mp3", "wav", "ogg", "m4a", "aac", "flac", "wma", "opus"}
    code_exts  = {"py", "js", "ts", "jsx", "tsx", "html", "css", "java", "c",
                  "cpp", "cs", "go", "rs", "rb", "php", "swift", "kt", "sh",
                  "bash", "ps1", "lua", "r", "m", "sql", "yaml", "toml"}
    archive_exts = {"zip", "rar", "tar", "gz", "7z", "bz2", "xz"}

    if ext in image_exts:   return "image"
    if ext in video_exts:   return "video"
    if ext in audio_exts:   return "audio"
    if ext in code_exts:    return "code"
    if ext in archive_exts: return "archive"
    if ext == "pdf":        return "pdf"
    if ext in ("docx", "doc"): return "docx"
    if ext in ("txt", "md", "rst", "log"): return "text"
    if ext in ("csv", "tsv"): return "csv"
    if ext in ("xlsx", "xls", "ods"): return "excel"
    if ext == "json":       return "json"
    if ext == "xml":        return "xml"
    if ext in ("pptx", "ppt"): return "pptx"
    return "unknown"


def _file_size_str(path: Path) -> str:
    size = path.stat().st_size
    if size < 1024:        return f"{size} B"
    if size < 1024**2:     return f"{size/1024:.1f} KB"
    if size < 1024**3:     return f"{size/1024**2:.1f} MB"
    return f"{size/1024**3:.1f} GB"

def _output_path(src: Path, suffix: str, new_ext: str = None) -> Path:
    ext  = new_ext or src.suffix
    name = f"{src.stem}_{suffix}{ext}"
    return src.parent / name

def _process_image(path: Path, action: str, params: dict, speak=None) -> str:
    try:
        from PIL import Image
    except ImportError:
        return "Pillow non installato. Esegui: pip install Pillow"

    action = action or "describe"

    if action in ("describe", "ocr", "analyze", "read", "extract_text"):
        try:
            model  = _gemini_client()
            img    = Image.open(path)
            prompt = {
                "describe": "Descrivi questa immagine in dettaglio in italiano.",
                "ocr":      "Estrai tutto il testo visibile nell'immagine.",
                "analyze":  "Analizza accuratamente questa immagine: oggetti, colori, contesto e testo.",
                "read":     "Leggi il testo contenuto preservandone la struttura.",
                "extract_text": "Estrai tutto il testo da questa immagine.",
            }.get(action, "Descrivi questa immagine.")

            if params.get("instruction"):
                prompt = params["instruction"]

            response = model.generate_content([prompt, img])
            result   = response.text.strip()

            if len(result) > 500 and params.get("save", True):
                out = _output_path(path, "risultato", ".txt")
                out.write_text(result, encoding="utf-8")
                return f"{result[:300]}...\n\nRisultato salvato in: {out}"
            return result
        except Exception as e:
            return f"Analisi immagine fallita: {e}"

    if action == "resize":
        width  = int(params.get("width",  0))
        height = int(params.get("height", 0))
        scale  = float(params.get("scale", 0))
        try:
            img = Image.open(path)
            w, h = img.size
            if scale:
                new_size = (int(w * scale), int(h * scale))
            elif width and height:
                new_size = (width, height)
            elif width:
                new_size = (width, int(h * width / w))
            elif height:
                new_size = (int(w * height / h), height)
            else:
                return "Specifica larghezza, altezza o scala."
            out = _output_path(path, f"ridimensionato_{new_size[0]}x{new_size[1]}")
            img.resize(new_size, Image.LANCZOS).save(out)
            return f"Ridimensionato da {w}x{h} a {new_size[0]}x{new_size[1]}. Salvato: {out.name}"
        except Exception as e:
            return f"Ridimensionamento fallito: {e}"

    if action == "convert":
        fmt = params.get("format", "png").lower().strip(".")
        fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "png": "PNG",
                   "webp": "WEBP", "bmp": "BMP", "tiff": "TIFF"}
        pil_fmt = fmt_map.get(fmt, fmt.upper())
        try:
            img = Image.open(path).convert("RGB") if fmt == "jpg" else Image.open(path)
            out = _output_path(path, "convertito", f".{fmt}")
            img.save(out, pil_fmt)
            return f"Convertito in {fmt.upper()}. Salvato: {out.name}"
        except Exception as e:
            return f"Conversione fallita: {e}"

    if action == "compress":
        quality = int(params.get("quality", 70))
        try:
            img = Image.open(path).convert("RGB")
            out = _output_path(path, f"compresso_q{quality}", ".jpg")
            img.save(out, "JPEG", quality=quality, optimize=True)
            before = _file_size_str(path)
            after  = _file_size_str(out)
            return f"Compresso: {before} → {after}. Salvato: {out.name}"
        except Exception as e:
            return f"Compressione fallita: {e}"

    if action == "info":
        try:
            img = Image.open(path)
            return (f"Info immagine: {img.format}, {img.size[0]}x{img.size[1]}px, "
                    f"modalità: {img.mode}, dimensione: {_file_size_str(path)}")
        except Exception as e:
            return f"Lettura info fallita: {e}"

    return _process_image(path, "describe", {"instruction": f"{action}: {params}"})

def _process_pdf(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"

    def _extract_pdf_text(max_chars=50000) -> str:
        text = ""
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
        except ImportError:
            try:
                import PyPDF2
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            except ImportError:
                return ""
        return text[:max_chars]

    if action in ("summarize", "extract_text", "translate_hint", "analyze", "reformat"):
        text = _extract_pdf_text()
        if not text.strip():
            return "Impossibile estrarre testo dal PDF (potrebbe essere scansionato/basato su immagini)."

        if action == "extract_text":
            out = _output_path(path, "testo_estratto", ".txt")
            out.write_text(text, encoding="utf-8")
            return f"Testo estratto ({len(text)} caratteri). Salvato: {out.name}"

        prompt_map = {
            "summarize":      f"Riassumi questo documento PDF in italiano in modo chiaro e strutturato:\n\n{text}",
            "analyze":        f"Analizza questo documento nel dettaglio:\n\n{text}",
            "translate_hint": f"Che lingua è e cosa dice il documento? Riassumi:\n\n{text}",
            "reformat":       f"Riformatta questo testo con una struttura pulita:\n\n{text}",
        }
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt_map.get(action, f"Analizza:\n\n{text}"))
            result   = response.text.strip()
            if len(result) > 600 and params.get("save", True):
                out = _output_path(path, action, ".txt")
                out.write_text(result, encoding="utf-8")
                return f"{result[:400]}...\n\nRisultato completo salvato: {out.name}"
            return result
        except Exception as e:
            return f"Analisi AI fallita: {e}"

    if action == "info":
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                pages = len(pdf.pages)
            return f"PDF: {pages} pagine, dimensione: {_file_size_str(path)}"
        except Exception:
            return f"Dimensione PDF: {_file_size_str(path)}"

    if action == "to_word":
        text = _extract_pdf_text()
        if not text:
            return "Impossibile estrarre il testo da convertire."
        try:
            from docx import Document
            doc  = Document()
            doc.add_heading(path.stem, 0)
            for para in text.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = _output_path(path, "convertito", ".docx")
            doc.save(out)
            return f"Convertito in documento Word. Salvato: {out.name}"
        except ImportError:
            return "python-docx non installato. Esegui: pip install python-docx"

    return f"Azione PDF non valida: '{action}'. Scegli tra: summarize, extract_text, info, to_word"

def _process_text_doc(path: Path, file_type: str, action: str,
                       params: dict, speak=None) -> str:
    action = action or "summarize"

    def _read_content() -> str:
        if file_type == "docx":
            try:
                from docx import Document
                doc  = Document(path)
                return "\n".join(p.text for p in doc.paragraphs)
            except ImportError:
                return "python-docx non installato."
            except Exception as e:
                return f"Lettura fallita: {e}"
        else:
            return path.read_text(encoding="utf-8", errors="ignore")

    content = _read_content()
    if not content.strip():
        return "Il file sembra vuoto."

    if action == "word_count":
        words = len(content.split())
        chars = len(content)
        lines = content.count("\n")
        return f"Conteggio: {words} parole, {chars} caratteri, {lines} righe."

    if action == "extract_text":
        if file_type != "txt":
            out = _output_path(path, "estratto", ".txt")
            out.write_text(content, encoding="utf-8")
            return f"Testo estratto. Salvato: {out.name}"
        return content[:2000]

    instruction = params.get("instruction", "")
    prompt_map  = {
        "summarize":  f"Riassumi questo documento in italiano:\n\n{content[:40000]}",
        "analyze":    f"Analizza questo documento:\n\n{content[:40000]}",
        "reformat":   f"Riformatta questo testo con titoli e paragrafi chiari:\n\n{content[:40000]}",
        "fix":        f"Correggi la grammatica e lo stile di questo testo:\n\n{content[:40000]}",
        "translate_hint": f"Che lingua è e cosa esprime? Riassumi:\n\n{content[:10000]}",
        "to_bullet":  f"Converti questo testo in punti elenco chiari:\n\n{content[:40000]}",
        "custom":     f"{instruction}\n\n{content[:40000]}",
    }

    if action not in prompt_map:
        action  = "custom"
        instruction = action

    try:
        model    = _gemini_client()
        response = model.generate_content(prompt_map[action])
        result   = response.text.strip()
        if len(result) > 600 and params.get("save", True):
            out = _output_path(path, action, ".txt")
            out.write_text(result, encoding="utf-8")
            return f"{result[:400]}...\n\nRisultato salvato: {out.name}"
        return result
    except Exception as e:
        return f"Elaborazione AI fallita: {e}"


def _process_data(path: Path, file_type: str, action: str,
                  params: dict, speak=None) -> str:
    try:
        import pandas as pd
    except ImportError:
        return "pandas non installato. Esegui: pip install pandas openpyxl"

    action = action or "analyze"

    try:
        if file_type == "csv":
            df = pd.read_csv(path, encoding="utf-8", errors="replace")
        else:
            df = pd.read_excel(path)
    except Exception as e:
        return f"Impossibile leggere il file dati: {e}"

    if action == "info":
        return (f"Righe: {len(df)}, Colonne: {len(df.columns)}\n"
                f"Colonne: {', '.join(df.columns.tolist())}\n"
                f"Dimensione: {_file_size_str(path)}")

    if action == "stats":
        try:
            desc = df.describe(include="all").to_string()
            return f"Statistiche:\n{desc[:2000]}"
        except Exception as e:
            return f"Statistiche fallite: {e}"

    if action == "analyze":
        preview = df.head(50).to_string()
        prompt  = (f"Analizza questo dataset. Colonne: {list(df.columns)}\n"
                   f"Righe: {len(df)}\nAnteprima:\n{preview}\n\n"
                   f"Fornisci spunti, pattern e dettagli rilevanti.")
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Analisi AI fallita: {e}"

    if action in ("convert", "to_csv", "to_excel", "to_json"):
        fmt = {"to_csv": "csv", "to_excel": "xlsx", "to_json": "json",
               "convert": params.get("format", "csv")}.get(action, "csv")
        try:
            if fmt == "csv":
                out = _output_path(path, "convertito", ".csv")
                df.to_csv(out, index=False, encoding="utf-8")
            elif fmt == "xlsx":
                out = _output_path(path, "convertito", ".xlsx")
                df.to_excel(out, index=False)
            elif fmt == "json":
                out = _output_path(path, "convertito", ".json")
                df.to_json(out, orient="records", force_ascii=False, indent=2)
            return f"Convertito in {fmt.UPPER()}. Salvato: {out.name}"
        except Exception as e:
            return f"Conversione fallita: {e}"

    if action == "filter":
        col       = params.get("column", "")
        value     = params.get("value", "")
        condition = params.get("condition", "equals")
        if not col or col not in df.columns:
            return f"Colonna '{col}' non trovata. Disponibili: {', '.join(df.columns)}"
        try:
            if condition == "equals":     filtered = df[df[col] == value]
            elif condition == "contains": filtered = df[df[col].astype(str).str.contains(str(value), case=False)]
            elif condition == "gt":       filtered = df[df[col] > float(value)]
            elif condition == "lt":       filtered = df[df[col] < float(value)]
            else:                         filtered = df[df[col] == value]
            out = _output_path(path, "filtrato", ".csv")
            filtered.to_csv(out, index=False)
            return f"Filtrato: {len(filtered)} righe trovate. Salvato: {out.name}"
        except Exception as e:
            return f"Filtro fallito: {e}"

    if action == "sort":
        col = params.get("column", df.columns[0])
        asc = params.get("ascending", True)
        try:
            sorted_df = df.sort_values(col, ascending=asc)
            out = _output_path(path, "ordinato", path.suffix)
            sorted_df.to_csv(out, index=False)
            return f"Ordinato per '{col}'. Salvato: {out.name}"
        except Exception as e:
            return f"Ordinamento fallito: {e}"

    preview = df.head(30).to_string()
    try:
        model    = _gemini_client()
        response = model.generate_content(
            f"Task: {action}\nDataset ({len(df)} righe, colonne: {list(df.columns)}):\n{preview}"
        )
        return response.text.strip()
    except Exception as e:
        return f"Elaborazione fallita: {e}"


def _process_json(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "analyze"
    try:
        content = path.read_text(encoding="utf-8")
        data    = json.loads(content)
    except Exception as e:
        return f"JSON non valido: {e}"

    if action == "validate":
        return f"JSON valido. Tipo: {type(data).__name__}, dimensione: {_file_size_str(path)}"

    if action == "format":
        out = _output_path(path, "formattato", ".json")
        out.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        return f"JSON formattato salvato: {out.name}"

    if action in ("analyze", "summarize", "extract"):
        preview = json.dumps(data, indent=2, ensure_ascii=False)[:8000]
        prompt  = f"Task: {action} questo JSON:\n{preview}"
        if params.get("instruction"):
            prompt = f"{params['instruction']}\n\nDati JSON:\n{preview}"
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Elaborazione AI fallita: {e}"

    if action == "to_csv":
        try:
            import pandas as pd
            if isinstance(data, list):
                df  = pd.DataFrame(data)
                out = _output_path(path, "convertito", ".csv")
                df.to_csv(out, index=False)
                return f"Convertito in CSV. Salvato: {out.name}"
            return "Il JSON deve essere una lista di oggetti per la conversione in CSV."
        except ImportError:
            return "pandas non installato."

    return _process_json(path, "analyze", {"instruction": action})

def _process_code(path: Path, action: str, params: dict, speak=None) -> str:
    action  = action or "explain"
    content = path.read_text(encoding="utf-8", errors="ignore")
    ext     = path.suffix.lstrip(".")

    if action == "run":
        if ext == "py":
            try:
                result = subprocess.run(
                    ["python", str(path)],
                    capture_output=True, text=True, timeout=30
                )
                out = result.stdout or result.stderr
                return f"Output:\n{out[:2000]}" if out else "Nessun output generato."
            except subprocess.TimeoutExpired:
                return "Esecuzione scaduta (timeout 30s)."
            except Exception as e:
                return f"Esecuzione fallita: {e}"
        return f"Esecuzione diretta non supportata per file .{ext}."

    if action == "info":
        lines = content.count("\n")
        words = len(content.split())
        return f"File di codice: {lines} righe, {words} parole, {_file_size_str(path)}"

    prompt_map = {
        "explain":   f"Spiega chiaramente questo codice {ext}:\n\n```{ext}\n{content[:30000]}\n```",
        "review":    f"Rivedi questo codice {ext} evidenziando bug, problemi e miglioramenti:\n\n```{ext}\n{content[:30000]}\n```",
        "fix":       f"Correggi eventuali bug in questo codice {ext} e restituisci la versione sistemata:\n\n```{ext}\n{content[:30000]}\n```",
        "optimize":  f"Ottimizza questo codice {ext} per prestazioni e leggibilità:\n\n```{ext}\n{content[:30000]}\n```",
        "document":  f"Aggiungi commenti e documentazione a questo codice {ext}:\n\n```{ext}\n{content[:30000]}\n```",
        "summarize": f"Riassumi la funzione di questo codice {ext}:\n\n```{ext}\n{content[:30000]}\n```",
        "test":      f"Scrivi unit test per questo codice {ext}:\n\n```{ext}\n{content[:30000]}\n```",
    }

    instruction = params.get("instruction", "")
    if action not in prompt_map:
        prompt = f"{action}\n\n```{ext}\n{content[:30000]}\n```"
        if instruction:
            prompt = f"{instruction}\n\n```{ext}\n{content[:30000]}\n```"
    else:
        prompt = prompt_map[action]

    try:
        model    = _gemini_client()
        response = model.generate_content(prompt)
        result   = response.text.strip()

        if action in ("fix", "optimize", "document") and params.get("save", True):
            out = _output_path(path, action)
            code_match = re.search(r"```(?:\w+)?\n(.*?)```", result, re.DOTALL)
            code_to_save = code_match.group(1) if code_match else result
            out.write_text(code_to_save, encoding="utf-8")
            return f"{result[:400]}...\n\nSalvato: {out.name}"
        return result
    except Exception as e:
        return f"Elaborazione AI fallita: {e}"

def _process_audio(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "transcribe"

    if action == "info":
        try:
            from pydub import AudioSegment
            audio    = AudioSegment.from_file(path)
            duration = len(audio) / 1000
            mins, secs = divmod(int(duration), 60)
            return (f"Audio: {mins}m {secs}s, "
                    f"{audio.channels} canali, "
                    f"{audio.frame_rate}Hz, "
                    f"{_file_size_str(path)}")
        except ImportError:
            return f"File Audio: {_file_size_str(path)} (installa pydub per maggiori dettagli)"
        except Exception as e:
            return f"Lettura info fallita: {e}"

    if action == "transcribe":
        try:
            model   = _gemini_client()
            content = path.read_bytes()
            mime    = {
                "mp3": "audio/mp3", "wav": "audio/wav",
                "ogg": "audio/ogg", "m4a": "audio/mp4",
                "aac": "audio/aac", "flac": "audio/flac",
            }.get(path.suffix.lstrip(".").lower(), "audio/mpeg")
            response = model.generate_content([
                "Trascrivi fedelmente tutto il parlato contenuto in questo file audio.",
                {"mime_type": mime, "data": content}
            ])
            result = response.text.strip()
            if params.get("save", True):
                out = _output_path(path, "trascrizione", ".txt")
                out.write_text(result, encoding="utf-8")
                return f"Trascrizione salvata: {out.name}\n\nAnteprima: {result[:300]}"
            return result
        except Exception as e:
            return f"Trascrizione fallita: {e}"

    if action == "convert":
        fmt = params.get("format", "mp3").lstrip(".")
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            out   = _output_path(path, "convertito", f".{fmt}")
            audio.export(out, format=fmt)
            return f"Convertito in {fmt.upper()}. Salvato: {out.name}"
        except ImportError:
            return "pydub non installato. Esegui: pip install pydub"
        except Exception as e:
            return f"Conversione fallita: {e}"

    if action == "trim":
        start = float(params.get("start", 0))
        end   = float(params.get("end",   0))
        try:
            from pydub import AudioSegment
            audio   = AudioSegment.from_file(path)
            end_ms  = int(end * 1000)   if end   else len(audio)
            trimmed = audio[int(start * 1000):end_ms]
            out     = _output_path(path, f"taglio_{int(start)}s_{int(end)}s")
            trimmed.export(out, format=path.suffix.lstrip("."))
            return f"Audio tagliato ({int(start)}s–{int(end)}s). Salvato: {out.name}"
        except ImportError:
            return "pydub non installato."
        except Exception as e:
            return f"Taglio fallito: {e}"

    return f"Azione audio sconosciuta: '{action}'. Prova: transcribe, info, convert, trim"

def _process_video(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "info"

    def _ffmpeg_available() -> bool:
        try:
            subprocess.run(["ffmpeg", "-version"], capture_output=True, timeout=3)
            return True
        except Exception:
            return False

    if action == "info":
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-print_format", "json",
                 "-show_format", "-show_streams", str(path)],
                capture_output=True, text=True, timeout=10
            )
            data     = json.loads(result.stdout)
            fmt      = data.get("format", {})
            duration = float(fmt.get("duration", 0))
            mins, secs = divmod(int(duration), 60)
            size     = _file_size_str(path)
            streams  = data.get("streams", [])
            video_s  = next((s for s in streams if s["codec_type"] == "video"), {})
            w        = video_s.get("width", "?")
            h        = video_s.get("height", "?")
            fps      = video_s.get("r_frame_rate", "?")
            return f"Video: {mins}m {secs}s, {w}x{h}, {fps} fps, {size}"
        except Exception:
            return f"File video: {_file_size_str(path)}"

    if action == "extract_audio":
        if not _ffmpeg_available():
            return "ffmpeg non trovato. Installa ffmpeg per estrarre l'audio."
        out = _output_path(path, "audio", ".mp3")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a", str(out), "-y"],
                capture_output=True, timeout=300
            )
            return f"Audio estratto. Salvato: {out.name}"
        except Exception as e:
            return f"Estrazione audio fallita: {e}"

    if action == "trim":
        start = params.get("start", "00:00:00")
        end   = params.get("end",   "")
        if not _ffmpeg_available():
            return "ffmpeg non trovato."
        out = _output_path(path, "taglio", path.suffix)
        try:
            cmd = ["ffmpeg", "-i", str(path), "-ss", str(start)]
            if end:
                cmd += ["-to", str(end)]
            cmd += ["-c", "copy", str(out), "-y"]
            subprocess.run(cmd, capture_output=True, timeout=600)
            return f"Video tagliato salvato: {out.name}"
        except Exception as e:
            return f"Taglio fallito: {e}"

    if action == "extract_frame":
        timestamp = params.get("timestamp", "00:00:01")
        if not _ffmpeg_available():
            return "ffmpeg non trovato."
        out = _output_path(path, f"frame_{timestamp.replace(':', '')}", ".jpg")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-ss", timestamp,
                 "-vframes", "1", str(out), "-y"],
                capture_output=True, timeout=30
            )
            return f"Frame estratto a {timestamp}. Salvato: {out.name}"
        except Exception as e:
            return f"Estrazione frame fallita: {e}"

    if action == "compress":
        crf = int(params.get("quality", 28))  
        if not _ffmpeg_available():
            return "ffmpeg non trovato."
        out = _output_path(path, f"compresso_crf{crf}", ".mp4")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path),
                 "-c:v", "libx264", "-crf", str(crf),
                 "-preset", "medium", "-c:a", "copy",
                 str(out), "-y"],
                capture_output=True, timeout=1800
            )
            before = _file_size_str(path)
            after  = _file_size_str(out)
            return f"Compresso: {before} → {after}. Salvato: {out.name}"
        except Exception as e:
            return f"Compressione fallita: {e}"

    if action == "transcribe":
        if not _ffmpeg_available():
            return "ffmpeg non trovato. Necessario per la trascrizione dei video."
        tmp_audio = Path(tempfile.mktemp(suffix=".mp3"))
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a",
                 str(tmp_audio), "-y"],
                capture_output=True, timeout=300
            )
            result = _process_audio(tmp_audio, "transcribe", params, speak)
            return result
        except Exception as e:
            return f"Trascrizione video fallita: {e}"
        finally:
            if tmp_audio.exists():
                tmp_audio.unlink()

    if action == "convert":
        fmt = params.get("format", "mp4").lstrip(".")
        if not _ffmpeg_available():
            return "ffmpeg non trovato."
        out = _output_path(path, "convertito", f".{fmt}")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), str(out), "-y"],
                capture_output=True, timeout=1800
            )
            return f"Convertito in {fmt.upper()}. Salvato: {out.name}"
        except Exception as e:
            return f"Conversione fallita: {e}"

    return f"Azione video sconosciuta: '{action}'. Prova: info, trim, extract_audio, extract_frame, compress, transcribe, convert"

def _process_archive(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "list"

    if action == "list":
        try:
            import zipfile, tarfile
            ext = path.suffix.lower()
            if ext == ".zip":
                with zipfile.ZipFile(path) as z:
                    names = z.namelist()
            elif ext in (".tar", ".gz", ".bz2", ".xz"):
                with tarfile.open(path) as t:
                    names = t.getnames()
            else:
                return f"Formato archivio non supportato: {ext}"
            preview = "\n".join(names[:30])
            suffix  = f"\n... e altri {len(names)-30}" if len(names) > 30 else ""
            return f"L'archivio contiene {len(names)} file:\n{preview}{suffix}"
        except Exception as e:
            return f"Lettura archivio fallita: {e}"

    if action == "extract":
        dest = Path(params.get("destination", str(path.parent / path.stem)))
        dest.mkdir(parents=True, exist_ok=True)
        try:
            shutil.unpack_archive(path, dest)
            return f"Estratto in: {dest}"
        except Exception as e:
            return f"Estrazione fallita: {e}"

    return f"Azione archivio sconosciuta: '{action}'. Prova: list, extract"

def _process_pptx(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"

    def _read_pptx_text() -> str:
        try:
            from pptx import Presentation
            prs  = Presentation(path)
            text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Diapositiva {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                text.append(slide_text)
            return "\n".join(text)
        except ImportError:
            return "python-pptx non installato."

    if action in ("summarize", "extract_text", "analyze"):
        text = _read_pptx_text()
        if action == "extract_text":
            out = _output_path(path, "testo", ".txt")
            out.write_text(text, encoding="utf-8")
            return f"Testo estratto. Salvato: {out.name}"
        try:
            model    = _gemini_client()
            prompt   = f"{'Riassumi' if action == 'summarize' else 'Analizza'} questa presentazione:\n{text[:30000]}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Elaborazione AI fallita: {e}"

    return f"Azione PPTX sconosciuta: '{action}'. Prova: summarize, extract_text, analyze"

def file_processor(parameters: dict, player=None, speak=None) -> str:
    file_path_str = parameters.get("file_path", "").strip()
    if not file_path_str:
        return "Nessun percorso file fornito."

    path = Path(file_path_str)
    if not path.exists():
        return f"File non trovato: {file_path_str}"
    if not path.is_file():
        return f"Il percorso non è un file: {file_path_str}"

    file_type   = _detect_type(path)
    action      = (parameters.get("action") or "").lower().strip()
    instruction = parameters.get("instruction", "")
    params      = {**parameters, "instruction": instruction}

    log_msg = f"[FileProcessor] {file_type.upper()} | {path.name} | action={action or 'auto'}"
    print(log_msg)
    if player:
        player.write_log(log_msg)

    if file_type == "unknown":
        try:
            content = path.read_text(encoding="utf-8", errors="ignore")[:10000]
            model   = _gemini_client()
            prompt  = f"File: {path.name}\nAnteprima contenuto:\n{content}\n\nTask: {action or instruction or 'Descrivi cosa contiene questo file.'}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Tipo file sconosciuto ({path.suffix}). Impossibile elaborare: {e}"

    dispatch = {
        "image":   _process_image,
        "pdf":     _process_pdf,
        "docx":    lambda p, a, pm, s: _process_text_doc(p, "docx", a, pm, s),
        "text":    lambda p, a, pm, s: _process_text_doc(p, "text", a, pm, s),
        "csv":     lambda p, a, pm, s: _process_data(p, "csv",   a, pm, s),
        "excel":   lambda p, a, pm, s: _process_data(p, "excel", a, pm, s),
        "json":    _process_json,
        "xml":     lambda p, a, pm, s: _process_json(p, a, pm, s),  
        "code":    _process_code,
        "audio":   _process_audio,
        "video":   _process_video,
        "archive": _process_archive,
        "pptx":    _process_pptx,
    }

    handler = dispatch.get(file_type)
    if not handler:
        return f"Tipo di file non supportato: {file_type}"

    try:
        result = handler(path, action, params, speak)
        return result or "Operazione completata."
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"Elaborazione fallita: {e}"